from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import numpy as np

# Create a presentation object
prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Pragyaan Rover Presentation"
subtitle.text = "An Overview of India's Lunar Rover"

# Slide 1 - Introduction
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Introduction to Pragyaan Rover"
content.text = (
    "Pragyaan, which means 'wisdom' in Sanskrit, is India's lunar rover designed for exploration "
    "of the Moon's surface. It was part of the Chandrayaan-2 mission launched by the Indian Space Research Organisation (ISRO)."
)

# Slide 2 - Specifications
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Specifications of Pragyaan Rover"
content.text = (
    "Weight: 27 kg\n"
    "Power: Solar powered\n"
    "Payloads: APXS (Alpha Particle X-ray Spectrometer) and LIBS (Laser Induced Breakdown Spectroscope)\n"
    "Speed: 1 cm per second\n"
    "Mission Duration: 14 Earth days (1 lunar day)\n"
    "Communication: With Vikram lander"
)

# Slide 3 - Mission Objectives
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Mission Objectives"
content.text = (
    "1. To study the lunar surface and gather data on its composition.\n"
    "2. To understand the presence of various elements like Magnesium, Aluminium, Silicon, Potassium, Calcium, Titanium, and Iron.\n"
    "3. To analyze the water ice in the lunar soil.\n"
    "4. To enhance the understanding of the Moon's origin and evolution."
)

# Slide 4 - Journey and Landing
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Journey and Landing"
content.text = (
    "Pragyaan was carried to the Moon aboard the Vikram lander of the Chandrayaan-2 mission.\n"
    "The landing site was selected near the lunar south pole, a region that holds immense scientific interest due to the presence of shadowed craters that may contain water ice."
)

# Slide 5 - Scientific Instruments
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Scientific Instruments"
content.text = (
    "1. Alpha Particle X-ray Spectrometer (APXS): Used for determining the elemental composition of the lunar soil.\n"
    "2. Laser Induced Breakdown Spectroscope (LIBS): Utilized for detecting and analyzing the presence of elements and minerals on the lunar surface."
)

# Slide 6 - Data and Findings
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Data and Findings"
content.text = (
    "The data collected by Pragyaan has contributed significantly to lunar science, offering insights into the composition and mineralogy of the Moon's surface.\n"
    "The findings help in understanding the geology of the Moon and assessing the presence of water ice in shadowed regions."
)

# Slide 7 - Conclusion
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Conclusion"
content.text = (
    "The Pragyaan rover marks a significant milestone in India's space exploration journey. Its successful deployment and data collection have paved the way for future missions and have contributed to global lunar science. Pragyaan has showcased India's growing capabilities in space technology and exploration."
)

# Save the presentation
prs.save('Pragyaan_Rover_Presentation.pptx')
print("Presentation saved successfully.")

